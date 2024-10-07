from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from tqdm import tqdm
import os
import logging

# Configure logging for this module
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

def create_table(slide, data, title):
    """Create a table on the given slide with the specified data and title.

    :param slide: The slide where the table will be added.
    :param data: The data for the table, provided as a list of lists.
    :param title: The title of the table.
    """
    rows, cols = len(data), len(data[0])
    left, top, width, height = Inches(1), Inches(1.5), Inches(11), Inches(4)

    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(24)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    for col_idx in range(cols):
        cell = table.cell(0, col_idx)
        cell.text = data[0][col_idx]
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 176, 240)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    for row_idx in range(1, rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(data[row_idx][col_idx])
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if row_idx % 2 == 0 else RGBColor(220, 230, 241)

def create_ppt(output_directory: str, presentation_path: str, template_path: str = None) -> None:
    """Create a PowerPoint presentation from analysis results.

    :param output_directory: The directory containing the analysis results.
    :param presentation_path: The path where the presentation will be saved.
    :param template_path: Optional path to a PowerPoint template.
    """
    prs = Presentation(template_path) if template_path else Presentation()
    if not template_path:
        prs.slide_width, prs.slide_height = Inches(13.3333), Inches(7.5)

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text, title_slide.placeholders[1].text = "EXPLORATORY DATA ANALYSIS", "MADE BY DORA"

    for paragraph in title_slide.shapes.title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.bold = True

    charts_files = os.listdir(f"{output_directory}/charts/")
    
    logging.info("Creating slides for charts...")
    for chart_file in tqdm(charts_files, desc="Processing Charts", ncols=100, unit="chart"):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.add_picture(f"{output_directory}/charts/{chart_file}", Inches(0.5), Inches(0.5), height=Inches(5))
        slide.shapes.title.text = f"Chart: {chart_file}"

    stats_file = os.path.join(output_directory, 'stats', 'univariate_analysis.txt')
    with open(stats_file, 'r') as f:
        stats_data = [line.split() for line in f.readlines()]

    create_table(prs.slides.add_slide(prs.slide_layouts[5]), stats_data, "Statistical Summary")

    prs.save(presentation_path)
    logging.info("Saved PowerPoint presentation successfully.")

def save_results(output_directory: str, template_path: str = None) -> None:
    """Save results of the analysis into a PowerPoint presentation.

    :param output_directory: The directory where analysis results are stored.
    :param template_path: Optional path to a PowerPoint template.
    """
    presentation_path = os.path.join(output_directory, 'EDA_Results.pptx')
    create_ppt(output_directory, presentation_path, template_path)
